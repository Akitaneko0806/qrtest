from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# プレゼンテーションの作成
prs = Presentation()

# スライドサイズの設定
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 色の設定
background_color = RGBColor(38, 38, 51)
yellow_color = RGBColor(255, 204, 0)
white_color = RGBColor(255, 255, 255)

# フォントサイズの設定
title_font_size = Pt(40)
subtitle_font_size = Pt(18)
small_font_size = Pt(12)

# スライド1: エグゼクティブ・サマリー
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = background_color

title_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "エグゼクティブ・サマリー"
title_frame.paragraphs[0].font.size = title_font_size
title_frame.paragraphs[0].font.color.rgb = white_color

# スライド2: 法令遵守
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = white_color

title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "法令遵守"
title_frame.paragraphs[0].font.size = title_font_size
title_frame.paragraphs[0].font.color.rgb = background_color

# スライド3: 組織図
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = white_color

title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "四門－組織図"
title_frame.paragraphs[0].font.size = title_font_size
title_frame.paragraphs[0].font.color.rgb = background_color

# テンプレートの保存
output_file = 'template.pptx'
prs.save(output_file)

print(f'Template saved as {output_file}')
